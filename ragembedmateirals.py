import os
import re
import json
import numpy as np
from collections import defaultdict, Counter
from pathlib import Path
import PyPDF2
from pptx import Presentation
from sentence_transformers import SentenceTransformer
from sklearn.metrics.pairwise import cosine_similarity
import networkx as nx

class TopicDocumentAssociator:
    """Associates predefined topics and concepts with documents and analyzes relationships"""
    
    def __init__(self, topics_file='topics.txt', concepts_file='concepts.txt', similarity_threshold=0.3):
        self.core_topics = self.load_items(topics_file, "topics")
        self.core_concepts = self.load_items(concepts_file, "concepts")
        self.similarity_threshold = similarity_threshold
        
        print("Loading models...")
        
        # Check for CUDA availability
        import torch
        self.device = 'cuda' if torch.cuda.is_available() else 'cpu'
        print(f"Using device: {self.device}")
        
        # Load sentence transformer for semantic matching
        self.sentence_model = SentenceTransformer('all-MiniLM-L6-v2', device=self.device)
        
        print(f"Loaded {len(self.core_topics)} core topics and {len(self.core_concepts)} concepts")
    
    def load_items(self, file_path, item_type):
        """Load topics or concepts from file"""
        items = []
        try:
            if os.path.exists(file_path):
                with open(file_path, 'r', encoding='utf-8') as f:
                    for line in f:
                        item = line.strip()
                        if item and len(item) > 2:
                            items.append(item)
                print(f"Loaded {len(items)} {item_type}")
            else:
                print(f"{item_type.title()} file '{file_path}' not found.")
                items = []
        except Exception as e:
            print(f"Could not load {item_type} file: {e}")
            items = []
        return items
    
    def extract_text_from_pdf(self, pdf_path):
        """Extract text from PDF file"""
        text = ""
        try:
            with open(pdf_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + " "
        except Exception as e:
            print(f"Error reading {pdf_path}: {e}")
        return text.strip()
    
    def extract_text_from_pptx(self, pptx_path):
        """Extract text from PowerPoint file"""
        text = ""
        try:
            prs = Presentation(pptx_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + " "
        except Exception as e:
            print(f"Error reading {pptx_path}: {e}")
        return text.strip()
    
    def find_exact_mentions(self, text, item):
        """Find exact mentions of a topic/concept in text - only true exact matches"""
        text_lower = text.lower()
        item_lower = item.lower()
        
        contexts = []
        total_count = 0
        
        # Strategy 1: Exact word boundary matches only
        pattern1 = r'\b' + re.escape(item_lower) + r'\b'
        matches1 = list(re.finditer(pattern1, text_lower))
        total_count += len(matches1)
        
        # Strategy 2: Allow for hyphenated versions
        item_hyphen = item_lower.replace(' ', '-')
        pattern2 = r'\b' + re.escape(item_hyphen) + r'\b'
        matches2 = list(re.finditer(pattern2, text_lower))
        total_count += len(matches2)
        
        # Strategy 3: Allow for plural forms (only if the base term was found exactly)
        if matches1 and not item_lower.endswith('s'):
            item_plural = item_lower + 's'
            pattern3 = r'\b' + re.escape(item_plural) + r'\b'
            matches3 = list(re.finditer(pattern3, text_lower))
            total_count += len(matches3)
        else:
            matches3 = []
        
        # NO partial word matching for multi-word terms - only exact complete matches
        
        # Extract contexts from all exact match types
        all_matches = matches1 + matches2 + matches3
        
        if all_matches:
            # Convert text to list for easier indexing
            words = text.split()
            text_positions = []
            current_pos = 0
            
            for word in words:
                text_positions.append(current_pos)
                current_pos += len(word) + 1  # +1 for space
            
            for match in all_matches:
                # Find which word index this match corresponds to
                match_start = match.start()
                word_idx = 0
                for i, pos in enumerate(text_positions):
                    if pos <= match_start < pos + len(words[i]):
                        word_idx = i
                        break
                
                # Extract context around this word
                start_idx = max(0, word_idx - 10)
                end_idx = min(len(words), word_idx + 15)
                context = ' '.join(words[start_idx:end_idx])
                
                contexts.append(context)
                
                if len(contexts) >= 5:  # Limit context examples
                    break
        
        return total_count, contexts
    
    def calculate_semantic_similarity(self, text, item):
        """Calculate semantic similarity with better relevance filtering"""
        # Split text into meaningful chunks (sentences/paragraphs)
        # Use both sentence and paragraph breaks for better context
        chunks = re.split(r'[.!?]+|\n\s*\n', text)
        chunks = [chunk.strip() for chunk in chunks if len(chunk.strip()) > 30]
        
        if not chunks:
            return 0.0, []
        
        # Filter out clearly irrelevant chunks (references, headers, etc.)
        filtered_chunks = []
        for chunk in chunks:
            chunk_lower = chunk.lower()
            # Skip reference sections, headers, etc.
            if not any(skip_word in chunk_lower for skip_word in 
                      ['references', 'bibliography', 'figure', 'table', 'appendix', 'doi:', 'http']):
                if len(chunk.split()) > 5:  # Ensure substantial content
                    filtered_chunks.append(chunk)
        
        if not filtered_chunks:
            return 0.0, []
        
        # Limit chunks for performance
        chunks_to_analyze = filtered_chunks[:30]
        
        # Get embeddings
        item_embedding = self.sentence_model.encode([item])
        chunk_embeddings = self.sentence_model.encode(chunks_to_analyze)
        
        # Calculate similarities
        similarities = cosine_similarity(chunk_embeddings, item_embedding).flatten()
        
        # Only keep chunks above a higher threshold for semantic relevance
        relevant_threshold = max(self.similarity_threshold, 0.4)
        relevant_indices = [i for i, sim in enumerate(similarities) if sim > relevant_threshold]
        
        if not relevant_indices:
            return 0.0, []
        
        # Get top relevant chunks
        relevant_similarities = [similarities[i] for i in relevant_indices]
        relevant_chunks = [chunks_to_analyze[i] for i in relevant_indices]
        
        # Sort by similarity and take top 3
        sorted_pairs = sorted(zip(relevant_similarities, relevant_chunks), reverse=True)
        top_similarities, top_chunks = zip(*sorted_pairs[:3])
        
        # Average similarity of relevant chunks only
        avg_similarity = np.mean(top_similarities)
        
        return float(avg_similarity), list(top_chunks)
    
    def analyze_document(self, file_path, doc_type, folder):
        """Analyze a single document for all topics and concepts"""
        print(f"  Analyzing {file_path.name}...")
        
        # Extract text
        if file_path.suffix.lower() == '.pdf':
            text = self.extract_text_from_pdf(file_path)
        else:
            text = self.extract_text_from_pptx(file_path)
        
        if not text or len(text) < 100:
            print(f"    Skipping {file_path.name} - insufficient text content")
            return None
        
        # Document info
        folder_name = Path(folder).name
        web_path = f"./{folder_name}/{file_path.name}"
        
        doc_info = {
            'name': file_path.name,
            'type': doc_type,
            'path': web_path,
            'text_length': len(text)
        }
        
        # Analyze all items (topics + concepts)
        item_analysis = {}
        
        # Analyze core topics
        for topic in self.core_topics:
            exact_count, exact_contexts = self.find_exact_mentions(text, topic)
            semantic_score, semantic_contexts = self.calculate_semantic_similarity(text, topic)
            
            # Calculate combined relevance score
            relevance = self.calculate_relevance_score(exact_count, semantic_score, len(text))
            
            # Only include if we have meaningful matches
            if exact_count > 0 or (semantic_score > 0.5 and len(semantic_contexts) > 0):
                item_analysis[topic] = {
                    'type': 'core_topic',
                    'exact_mentions': exact_count,
                    'semantic_score': semantic_score,
                    'relevance_score': relevance,
                    'exact_contexts': exact_contexts,
                    'semantic_contexts': semantic_contexts,
                    'color': '#9333ea'  # Purple for core topics
                }
        
        # Analyze concepts
        for concept in self.core_concepts:
            exact_count, exact_contexts = self.find_exact_mentions(text, concept)
            semantic_score, semantic_contexts = self.calculate_semantic_similarity(text, concept)
            
            # Calculate combined relevance score
            relevance = self.calculate_relevance_score(exact_count, semantic_score, len(text))
            
            if exact_count > 0 or semantic_score > self.similarity_threshold:
                item_analysis[concept] = {
                    'type': 'concept',
                    'exact_mentions': exact_count,
                    'semantic_score': semantic_score,
                    'relevance_score': relevance,
                    'exact_contexts': exact_contexts,
                    'semantic_contexts': semantic_contexts,
                    'color': '#059669'  # Green for concepts
                }
        
        return doc_info, item_analysis
    
    def calculate_relevance_score(self, exact_count, semantic_score, text_length):
        """Calculate combined relevance score with heavy preference for exact matches"""
        # Normalize exact mentions by document length (per 1000 words)
        exact_normalized = (exact_count / max(text_length / 1000, 1)) * 10
        
        # Give much higher weight to exact matches
        if exact_count > 0:
            # If we have exact matches, weight them very heavily
            combined_score = (exact_normalized * 0.9) + (semantic_score * 0.1)
            # Boost for multiple exact mentions
            if exact_count > 1:
                combined_score *= 1.2
        else:
            # Only semantic similarity, but with lower base score
            combined_score = semantic_score * 0.5
        
        return min(combined_score, 10.0)  # Cap at 10
    
    def process_folders(self, slides_folder, literature_folder):
        """Process all documents and create associations"""
        print("Processing documents for topic and concept associations...")
        
        documents = []
        all_associations = defaultdict(lambda: {
            'docs': {},
            'type': '',
            'total_mentions': 0,
            'total_semantic_score': 0.0,
            'avg_relevance': 0.0,
            'color': '#4c9aff'
        })
        
        # Collect all files
        all_files = []
        for folder, doc_type in [(slides_folder, 'slide'), (literature_folder, 'literature')]:
            if os.path.exists(folder):
                for file_path in Path(folder).glob('*'):
                    if file_path.suffix.lower() in ['.pdf', '.pptx']:
                        all_files.append((file_path, doc_type, folder))
        
        if not all_files:
            print("No PDF or PPTX files found in the specified folders.")
            return [], {}
        
        # Process each file
        for file_path, doc_type, folder in all_files:
            result = self.analyze_document(file_path, doc_type, folder)
            
            if result is None:
                continue
                
            doc_info, item_analysis = result
            documents.append(doc_info)
            
            # Aggregate associations
            for item_name, analysis in item_analysis.items():
                # Update document-specific data
                all_associations[item_name]['docs'][doc_info['name']] = {
                    'exact_mentions': analysis['exact_mentions'],
                    'semantic_score': analysis['semantic_score'],
                    'relevance_score': analysis['relevance_score'],
                    'exact_contexts': analysis['exact_contexts'],
                    'semantic_contexts': analysis['semantic_contexts']
                }
                
                # Update aggregate data
                all_associations[item_name]['type'] = analysis['type']
                all_associations[item_name]['color'] = analysis['color']
                all_associations[item_name]['total_mentions'] += analysis['exact_mentions']
                all_associations[item_name]['total_semantic_score'] += analysis['semantic_score']
        
        # Calculate averages and prepare final data
        final_associations = {}
        for item_name, data in all_associations.items():
            if len(data['docs']) > 0:  # Only include items found in documents
                # Calculate average relevance
                relevances = [doc_data['relevance_score'] for doc_data in data['docs'].values()]
                avg_relevance = sum(relevances) / len(relevances)
                
                # Sort document names by exact matches first, then relevance
                sorted_doc_names = sorted(data['docs'].keys(), key=lambda doc_name: (
                    -data['docs'][doc_name]['exact_mentions'],  # Negative for descending
                    -data['docs'][doc_name]['relevance_score']   # Negative for descending
                ))
                
                final_associations[item_name] = {
                    'docs': sorted_doc_names,  # Now properly sorted
                    'doc_details': data['docs'],
                    'type': data['type'],
                    'total_mentions': data['total_mentions'],
                    'avg_semantic_score': data['total_semantic_score'] / len(data['docs']),
                    'avg_relevance': avg_relevance,
                    'document_count': len(data['docs']),
                    'color': data['color']
                }
        
        # Print summary
        print(f"\nFound associations for {len(final_associations)} items:")
        topics_found = sum(1 for item in final_associations.values() if item['type'] == 'core_topic')
        concepts_found = sum(1 for item in final_associations.values() if item['type'] == 'concept')
        print(f"  - {topics_found} core topics")
        print(f"  - {concepts_found} concepts")
        print(f"  - Across {len(documents)} documents")
        
        return documents, final_associations
    
    def analyze_relationships(self, associations):
        """Analyze relationships between topics and concepts based on document co-occurrence"""
        relationships = []
        
        items = list(associations.keys())
        for i, item1 in enumerate(items):
            docs1 = set(associations[item1]['docs'])
            
            for item2 in items[i+1:]:
                docs2 = set(associations[item2]['docs'])
                shared_docs = docs1.intersection(docs2)
                
                if len(shared_docs) > 0:
                    # Calculate relationship strength
                    jaccard_similarity = len(shared_docs) / len(docs1.union(docs2))
                    overlap_ratio = len(shared_docs) / min(len(docs1), len(docs2))
                    
                    # Combined strength score
                    strength = (jaccard_similarity * 0.6) + (overlap_ratio * 0.4)
                    
                    relationships.append({
                        'item1': item1,
                        'item2': item2,
                        'shared_documents': len(shared_docs),
                        'shared_doc_names': list(shared_docs),
                        'strength': strength,
                        'jaccard_similarity': jaccard_similarity,
                        'overlap_ratio': overlap_ratio
                    })
        
        # Sort by strength
        relationships.sort(key=lambda x: x['strength'], reverse=True)
        return relationships


def create_topic_concept_network(documents, associations, relationships, output_file='topic_concept_network.html'):
    """Create interactive network visualization"""
    
    G = nx.Graph()
    
    # Add nodes for topics and concepts
    for item_name, data in associations.items():
        doc_count = data['document_count']
        total_mentions = data['total_mentions']
        avg_relevance = data['avg_relevance']
        item_type = data['type']
        
        # Calculate node size based on relevance and document coverage (tiny, consistent sizes)
        if item_type == 'core_topic':
            base_size = 4
            size = base_size + (avg_relevance * 0.3) + (doc_count * 0.2)
            font_size = 7 + int(avg_relevance * 0.1) + (doc_count * 0.05)
        else:  # concept
            base_size = 3
            size = base_size + (avg_relevance * 0.2) + (doc_count * 0.15)
            font_size = 6 + int(avg_relevance * 0.1) + (doc_count * 0.05)
        
        size = min(size, 6)  # Very small cap
        font_size = min(font_size, 7)  # Very small font cap
        
        G.add_node(item_name,
                  size=int(size),
                  font_size=int(font_size),
                  doc_count=doc_count,
                  total_mentions=total_mentions,
                  avg_relevance=round(avg_relevance, 2),
                  item_type=item_type,
                  color=data['color'])
    
    # Add edges based on relationships with better weighting
    for rel in relationships:
        if rel['strength'] > 0.05:  # Lower threshold to show more connections
            # Calculate edge weight based on relationship strength
            edge_weight = max(1, rel['strength'] * 15)  # Scale for visibility
            edge_width = max(1, rel['strength'] * 8)    # Visual width
            
            G.add_edge(rel['item1'], rel['item2'],
                      weight=rel['strength'],
                      shared_docs=rel['shared_documents'],
                      shared_doc_names=rel['shared_doc_names'],
                      strength=round(rel['strength'], 3),
                      edge_weight=edge_weight,
                      edge_width=edge_width)
    
    # Convert to vis.js format
    nodes = []
    for node, attrs in G.nodes(data=True):
        nodes.append({
            'id': node,
            'label': node.title(),
            'value': attrs['size'],
            'title': f"{node}<br>Type: {attrs['item_type']}<br>Documents: {attrs['doc_count']}<br>Avg Relevance: {attrs['avg_relevance']}<br>Total Mentions: {attrs['total_mentions']}",
            'color': attrs['color'],
            'font': {'size': attrs['font_size']},
            'type': attrs['item_type'],
            'doc_count': attrs['doc_count'],
            'avg_relevance': attrs['avg_relevance'],
            'size': attrs['size']  # Explicitly set size to ensure it's used
        })
    
    edges = []
    for u, v, attrs in G.edges(data=True):
        # Use very minimal edge weight and width
        edge_width = 0.1  # Extremely thin
        edge_weight = 1   # Minimal weight
        
        # Color edges based on strength
        edge_color = '#e2e8f0'  # Very light gray
        if attrs['strength'] > 0.3:
            edge_color = '#cbd5e0'  # slightly darker for strong connections
        
        edges.append({
            'from': u,
            'to': v,
            'value': edge_weight,
            'width': edge_width,
            'title': f"Relationship Strength: {attrs['strength']}<br>Shared Documents: {attrs['shared_docs']}<br>Documents: {', '.join(attrs['shared_doc_names'])}",
            'color': {
                'color': edge_color,
                'highlight': '#a0aec0',
                'hover': '#a0aec0'
            },
            'strength': attrs['strength']
        })
    
    # Write HTML file using string concatenation instead of f-strings
    html_parts = []
    
    # HTML header
    html_parts.append("""<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <title>Topic & Concept Network</title>
    <script src="https://unpkg.com/vis-network@9.1.2/standalone/umd/vis-network.min.js"></script>
    <style>
        * { margin: 0; padding: 0; box-sizing: border-box; }
        body { font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif; }
        
        #controls {
            position: absolute;
            top: 20px;
            left: 20px;
            z-index: 1000;
            background: white;
            padding: 15px;
            border-radius: 8px;
            box-shadow: 0 2px 8px rgba(0,0,0,0.1);
        }
        
        #search-input {
            padding: 10px 14px;
            border: 2px solid #e0e0e0;
            border-radius: 4px;
            width: 320px;
            font-size: 16px;
            margin-bottom: 10px;
        }
        
        .filter-btn {
            padding: 6px 12px;
            margin: 2px;
            border: 1px solid #4c9aff;
            border-radius: 4px;
            background: white;
            color: #4c9aff;
            cursor: pointer;
            font-size: 12px;
        }
        
        .filter-btn.active {
            background: #4c9aff;
            color: white;
        }
        
        #reset-btn {
            padding: 8px 16px;
            margin-left: 10px;
            border: 1px solid #4c9aff;
            border-radius: 4px;
            background: #4c9aff;
            color: white;
            cursor: pointer;
            font-size: 14px;
        }
        
        #legend {
            position: absolute;
            bottom: 20px;
            left: 20px;
            z-index: 1000;
            background: white;
            padding: 12px 15px;
            border-radius: 8px;
            box-shadow: 0 2px 8px rgba(0,0,0,0.1);
            font-size: 13px;
        }
        
        .legend-item {
            display: flex;
            align-items: center;
            margin: 5px 0;
        }
        
        .legend-color {
            width: 20px;
            height: 20px;
            border-radius: 50%;
            margin-right: 8px;
            border: 2px solid #e0e0e0;
        }
        
        #network { width: 65%; height: 100vh; float: left; background: #fafafa; }
        #panel { 
            width: 35%; 
            height: 100vh; 
            float: right; 
            background: white;
            overflow-y: auto;
            padding: 20px;
            border-left: 2px solid #e0e0e0;
        }
        
        .item-title { 
            font-size: 24px; 
            font-weight: 700; 
            margin-bottom: 15px;
            color: #1a202c;
            border-bottom: 3px solid #4c9aff;
            padding-bottom: 8px;
        }
        
        .item-info {
            background: #f0f9ff;
            padding: 12px;
            border-radius: 6px;
            margin-bottom: 15px;
            font-size: 14px;
        }
        
        .doc-section {
            margin-bottom: 20px;
        }
        
        .doc-item { 
            background: #f7fafc;
            padding: 12px;
            margin: 8px 0;
            border-radius: 6px;
            cursor: pointer;
            border-left: 4px solid #4c9aff;
            transition: all 0.2s;
            font-size: 14px;
        }
        
        .doc-item:hover { 
            background: #edf2f7;
            transform: translateX(4px);
        }
        
        .doc-name {
            font-weight: 600;
            margin-bottom: 5px;
        }
        
        .doc-stats {
            color: #718096;
            font-size: 12px;
            margin-bottom: 8px;
        }
        
        .context-snippet {
            background: #fff;
            padding: 8px;
            border-radius: 4px;
            margin: 4px 0;
            font-style: italic;
            font-size: 12px;
            border-left: 3px solid #e2e8f0;
        }
        
        #debug {
            position: absolute;
            top: 20px;
            right: 20px;
            background: #fffbeb;
            padding: 10px;
            border-radius: 4px;
            font-size: 12px;
            border: 1px solid #fbbf24;
        }
    </style>
</head>
<body>
    <div id="debug">
        <div>Nodes: <span id="node-count">Loading...</span></div>
        <div>Edges: <span id="edge-count">Loading...</span></div>
        <div>Status: <span id="status">Initializing...</span></div>
    </div>
    
    <div id="controls">
        <input type="text" id="search-input" placeholder="Search topics and concepts..." />
        <button id="reset-btn" onclick="resetView()">Reset View</button>
        <br>
        <button class="filter-btn active" onclick="filterNodes('all')">All</button>
        <button class="filter-btn" onclick="filterNodes('core_topic')">Core Topics</button>
        <button class="filter-btn" onclick="filterNodes('concept')">Concepts</button>
        <br>
        <button class="filter-btn" onclick="fitNetwork()">Fit to View</button>
    </div>
    
    <div id="legend">
        <div class="legend-item">
            <div class="legend-color" style="background: #9333ea;"></div>
            <span>Core Topics</span>
        </div>
        <div class="legend-item">
            <div class="legend-color" style="background: #059669;"></div>
            <span>Concepts</span>
        </div>
    </div>
    
    <div id="network"></div>
    <div id="panel">
        <div style="text-align: center; color: #a0aec0; padding: 60px 20px;">
            Click a topic or concept to explore its document associations<br>
            <small>Node size indicates relevance and document coverage</small>
        </div>
    </div>
    
    <script>
        // Debug function
        function updateDebug(status) {
            document.getElementById('status').textContent = status;
            console.log('Status:', status);
        }
        
        // Initialize data
        updateDebug('Loading data...');
        
        const nodesData = """)
    
    # Add nodes data
    html_parts.append(json.dumps(nodes, indent=2))
    html_parts.append(";\n        const edgesData = ")
    html_parts.append(json.dumps(edges, indent=2))
    html_parts.append(";\n        const allAssociations = ")
    
    # Prepare associations data
    associations_data = {
        k: {
            **v,
            'doc_details': v.get('doc_details', {})
        } for k, v in associations.items()
    }
    html_parts.append(json.dumps(associations_data, indent=2))
    html_parts.append(";\n        const docPaths = ")
    html_parts.append(json.dumps({d['name']: d['path'] for d in documents}, indent=2))
    
    # Continue with JavaScript
    html_parts.append(""";
        
        updateDebug('Data loaded, creating network...');
        
        // Update debug counters
        document.getElementById('node-count').textContent = nodesData.length;
        document.getElementById('edge-count').textContent = edgesData.length;
        
        // Create vis.js datasets
        const nodes = new vis.DataSet(nodesData);
        const edges = new vis.DataSet(edgesData);
        
        const container = document.getElementById('network');
        const data = { nodes: nodes, edges: edges };
        
        const options = {
            physics: {
                enabled: true,
                stabilization: {
                    enabled: true,
                    iterations: 2000,
                    updateInterval: 100,
                    onlyDynamicEdges: false,
                    fit: true
                },
                barnesHut: {
                    gravitationalConstant: -2000,
                    centralGravity: 0.3,
                    springLength: 95,
                    springConstant: 0.04,
                    damping: 0.95,
                    avoidOverlap: 0.1
                },
                maxVelocity: 50,
                minVelocity: 0.1,
                solver: 'barnesHut',
                timestep: 0.5,
                adaptiveTimestep: true
            },
            layout: {
                randomSeed: 42,
                improvedLayout: true
            },
            interaction: { 
                hover: true,
                hoverConnectedEdges: true,
                dragNodes: true,
                dragView: true,
                zoomView: true
            },
            nodes: {
                shape: 'dot',
                font: { 
                    color: '#1a202c',
                    face: 'Arial',
                    size: 8,
                    strokeWidth: 1,
                    strokeColor: '#ffffff',
                    vadjust: 0,
                    align: 'center'
                },
                borderWidth: 1,
                shadow: false,
                scaling: {
                    min: 3,
                    max: 6
                }
            },
            edges: {
                smooth: { 
                    type: 'continuous',
                    forceDirection: 'none',
                    roundness: 0.3
                },
                color: '#cbd5e0',
                width: 0.1,
                hoverWidth: 0.3,
                scaling: {
                    min: 0.1,
                    max: 0.1
                }
            }
        };
        
        updateDebug('Creating network visualization...');
        
        const network = new vis.Network(container, data, options);
        
        network.on('stabilizationIterationsDone', function() {
            network.setOptions({ physics: false });
            updateDebug('Network stabilized and ready!');
        });
        
        network.on('click', function(params) {
            if (params.nodes.length > 0) {
                showItemDetails(params.nodes[0]);
            }
        });
        
        updateDebug('Network created, stabilizing...');
        
        function stopPhysics() {
            network.setOptions({ physics: false });
            updateDebug('Physics stopped - network is now static');
        }
        
        function fitNetwork() {
            network.fit({
                animation: {
                    duration: 1000,
                    easingFunction: 'easeInOutQuad'
                }
            });
        }
        
        function showItemDetails(itemName) {
            const data = allAssociations[itemName];
            if (!data) {
                console.error('No data found for item:', itemName);
                return;
            }
            
            let html = '<div class="item-title">' + itemName + '</div>';
            
            // Item information
            html += '<div class="item-info">';
            html += '<strong>Type:</strong> ' + data.type.replace('_', ' ').toUpperCase() + '<br>';
            html += '<strong>Documents:</strong> ' + data.document_count + '<br>';
            html += '<strong>Total Mentions:</strong> ' + data.total_mentions + '<br>';
            html += '<strong>Avg Relevance:</strong> ' + data.avg_relevance.toFixed(2) + '<br>';
            html += '<strong>Avg Semantic Score:</strong> ' + data.avg_semantic_score.toFixed(3);
            html += '</div>';
            
            // Document details
            html += '<div class="doc-section">';
            html += '<div style="color: #718096; font-weight: 600; margin: 15px 0 10px 0;">Document Associations (sorted by relevance):</div>';
            
            // Sort documents by exact matches first, then by relevance score
            const sortedDocs = data.docs.slice().sort(function(a, b) {
                const detailsA = data.doc_details[a];
                const detailsB = data.doc_details[b];
                
                // First priority: exact matches
                if (detailsA.exact_mentions !== detailsB.exact_mentions) {
                    return detailsB.exact_mentions - detailsA.exact_mentions;
                }
                
                // Second priority: relevance score
                return detailsB.relevance_score - detailsA.relevance_score;
            });
            
            sortedDocs.forEach(function(docName) {
                const docDetails = data.doc_details[docName];
                const path = docPaths[docName];
                
                html += '<div class="doc-item" onclick="window.open(\\'' + path + '\\', \\'_blank\\')">';
                html += '<div class="doc-name">' + docName + '</div>';
                html += '<div class="doc-stats">';
                html += 'Exact Mentions: ' + docDetails.exact_mentions + ' | ';
                html += 'Semantic Score: ' + docDetails.semantic_score.toFixed(3) + ' | ';
                html += 'Relevance: ' + docDetails.relevance_score.toFixed(2);
                html += '</div>';
                
                // Add context snippets
                if (docDetails.exact_contexts && docDetails.exact_contexts.length > 0) {
                    html += '<div style="margin-top: 8px;"><strong>Exact mentions:</strong>';
                    docDetails.exact_contexts.slice(0, 2).forEach(function(context) {
                        html += '<div class="context-snippet">"...' + context + '..."</div>';
                    });
                    html += '</div>';
                }
                
                if (docDetails.semantic_contexts && docDetails.semantic_contexts.length > 0) {
                    html += '<div style="margin-top: 8px;"><strong>Related content:</strong>';
                    docDetails.semantic_contexts.slice(0, 2).forEach(function(context) {
                        html += '<div class="context-snippet">"' + context.substring(0, 150) + '..."</div>';
                    });
                    html += '</div>';
                }
                
                html += '</div>';
            });
            
            html += '</div>';
            
            document.getElementById('panel').innerHTML = html;
        }
        
        function filterNodes(type) {
            // Update button states
            document.querySelectorAll('.filter-btn').forEach(function(btn) { 
                btn.classList.remove('active'); 
            });
            event.target.classList.add('active');
            
            // Filter nodes
            const allNodes = nodes.get();
            const filteredNodes = type === 'all' ? allNodes : allNodes.filter(function(node) { 
                return node.type === type; 
            });
            
            // Get edges connected to visible nodes
            const visibleNodeIds = new Set(filteredNodes.map(function(n) { return n.id; }));
            const allEdges = edges.get();
            const filteredEdges = allEdges.filter(function(edge) {
                return visibleNodeIds.has(edge.from) && visibleNodeIds.has(edge.to);
            });
            
            // Update network
            network.setData({ nodes: filteredNodes, edges: filteredEdges });
            network.fit();
        }
        
        function resetView() {
            // Reset all data
            network.setData({ nodes: nodes.get(), edges: edges.get() });
            
            // Reset search
            document.getElementById('search-input').value = '';
            
            // Reset filter
            document.querySelectorAll('.filter-btn').forEach(function(btn) { 
                btn.classList.remove('active'); 
            });
            document.querySelector('.filter-btn').classList.add('active');
            
            network.fit();
        }
        
        // Search functionality
        document.getElementById('search-input').addEventListener('input', function(e) {
            const searchTerm = e.target.value.toLowerCase().trim();
            const allNodes = nodes.get();
            
            if (searchTerm.length > 2) {
                const matches = allNodes.filter(function(node) {
                    return node.id.toLowerCase().includes(searchTerm);
                });
                
                if (matches.length > 0) {
                    // Highlight matches
                    const updatedNodes = allNodes.map(function(node) {
                        return {
                            ...node,
                            color: node.id.toLowerCase().includes(searchTerm) ? '#f59e0b' : allAssociations[node.id] ? allAssociations[node.id].color : '#4c9aff'
                        };
                    });
                    
                    nodes.update(updatedNodes);
                    
                    // Focus on best match
                    matches.sort(function(a, b) { 
                        return b.avg_relevance - a.avg_relevance; 
                    });
                    network.focus(matches[0].id, { scale: 1.2, animation: true });
                    showItemDetails(matches[0].id);
                }
            } else {
                // Reset colors
                const updatedNodes = allNodes.map(function(node) {
                    return {
                        ...node,
                        color: allAssociations[node.id] ? allAssociations[node.id].color : '#4c9aff'
                    };
                });
                nodes.update(updatedNodes);
            }
        });
        
        console.log('Loaded network with ' + Object.keys(allAssociations).length + ' items');
        console.log('Nodes:', nodesData.length, 'Edges:', edgesData.length);
        
        // Test if data is properly loaded
        if (nodesData.length === 0) {
            document.getElementById('panel').innerHTML = '<div style="padding: 20px; color: red;">No nodes found! Check if your topics/concepts files exist and contain data.</div>';
        }
        
        // Hide debug panel after initialization (optional)
        setTimeout(function() {
            document.getElementById('debug').style.display = 'none';
        }, 5000);  // Hide after 5 seconds
    </script>
</body>
</html>""")
    
    # Join all HTML parts
    html_content = ''.join(html_parts)
    
    with open(output_file, 'w', encoding='utf-8') as f:
        f.write(html_content)
    
    print(f"\nCreated {output_file}")
    print(f"Network contains {len(nodes)} nodes and {len(edges)} edges")
    print("Features:")
    print("  - Core topics (purple) and concepts (green) from separate files")
    print("  - Document associations with relevance scoring")
    print("  - Relationship analysis based on document co-occurrence")
    print("  - Interactive filtering and search")
    print("  - Context snippets showing actual usage")
    print("  - Debug panel to help diagnose issues")


def main():
    """
    Main function to process topics, concepts, and documents
    """
    slides_folder = './pdfs'
    literature_folder = './pdfs'
    topics_file = 'topics.txt'
    concepts_file = 'concepts.txt'
    
    print("Topic and Concept Document Association Tool")
    print("=" * 50)
    
    # Check if folders exist
    if not os.path.exists(slides_folder) and not os.path.exists(literature_folder):
        print(f"Warning: Neither folder exists - {slides_folder} or {literature_folder}")
        print("Please ensure your documents are in the correct folders.")
        return
    
    # Check if topic/concept files exist
    if not os.path.exists(topics_file):
        print(f"Warning: {topics_file} not found. Creating sample file...")
        with open(topics_file, 'w', encoding='utf-8') as f:
            f.write("consciousness\ncognitive science\nneural networks\nmachine learning\nartificial intelligence\n")
        print(f"Created sample {topics_file}. Please edit it with your core topics.")
    
    if not os.path.exists(concepts_file):
        print(f"Warning: {concepts_file} not found. Creating sample file...")
        with open(concepts_file, 'w', encoding='utf-8') as f:
            f.write("attention\nmemory\nperception\nlearning\ndecision making\nemotion\n")
        print(f"Created sample {concepts_file}. Please edit it with your concepts.")
    
    # Initialize associator
    associator = TopicDocumentAssociator(topics_file, concepts_file, similarity_threshold=0.3)
    
    # Process documents
    documents, associations = associator.process_folders(slides_folder, literature_folder)
    
    if associations:
        # Analyze relationships
        relationships = associator.analyze_relationships(associations)
        
        # Create network visualization
        create_topic_concept_network(documents, associations, relationships, 'topic_concept_network.html')
        
        # Print summary statistics
        print(f"\nSummary Statistics:")
        print(f"  - {len(documents)} documents processed")
        print(f"  - {len(associations)} topics/concepts found in documents")
        print(f"  - {len(relationships)} relationships identified")
        
        # Show top associations
        sorted_associations = sorted(associations.items(), 
                                   key=lambda x: (x[1]['document_count'], x[1]['avg_relevance']), 
                                   reverse=True)
        
        print(f"\nTop 10 associations by document coverage:")
        for i, (name, data) in enumerate(sorted_associations[:10]):
            print(f"  {i+1:2d}. {name} ({data['type']}) - {data['document_count']} docs, {data['avg_relevance']:.2f} avg relevance")
        
        # Show top relationships
        if relationships:
            print(f"\nTop 5 strongest relationships:")
            for i, rel in enumerate(relationships[:5]):
                print(f"  {i+1}. {rel['item1']} ↔ {rel['item2']} (strength: {rel['strength']:.3f}, {rel['shared_documents']} shared docs)")
        
        print(f"\nOpen 'topic_concept_network.html' in your browser to explore the interactive network!")
        
    else:
        print("No associations found. Check if:")
        print("  1. Documents contain readable text")
        print("  2. Topics and concepts in your .txt files match content in documents")
        print("  3. Similarity threshold is appropriate (currently 0.3)")


if __name__ == "__main__":
    main()